import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

prs = Presentation(r'd:\laragon\www\andrian\Presentasi_Sistem_Kredit_BPR_Wonosobo.pptx')
slide5 = prs.slides[4]

# ============================================
# STEP 1: Hapus SEMUA shape kecuali background
# ============================================
# Keep only Shape 0 (number badge "2") - we'll rebuild everything
shapes_to_keep = set()
for i, shape in enumerate(slide5.shapes):
    # Keep the dark sidebar if exists (type AUTO_SHAPE at left=0)
    pass

# Actually, let's delete ALL shapes and rebuild from scratch for a clean slate
while len(slide5.shapes) > 0:
    sp = slide5.shapes[0]._element
    sp.getparent().remove(sp)

# ============================================
# CONSTANTS
# ============================================
SLIDE_W = 12191695  # 13.33 inches
SLIDE_H = 6858000   # 7.50 inches

dark_blue = RGBColor(29, 56, 88)
gold = RGBColor(196, 154, 69)
white = RGBColor(255, 255, 255)
light_bg = RGBColor(240, 244, 248)
red_text = RGBColor(192, 57, 43)
grey_text = RGBColor(100, 100, 100)

# Jenis kredit card colors
card_colors = [
    RGBColor(41, 128, 185),   # Umum - Blue
    RGBColor(39, 174, 96),    # PPPK - Green
    RGBColor(142, 68, 173),   # Perangkat Desa - Purple
    RGBColor(230, 126, 34),   # KPR - Orange
    RGBColor(241, 196, 15),   # Kretamas - Yellow
    RGBColor(44, 62, 80),     # Cash Collateral - Dark
]

# Status colors
status_colors = [
    RGBColor(149, 165, 166),  # DRAFT - Grey
    RGBColor(41, 128, 185),   # DIAJUKAN - Blue
    RGBColor(230, 126, 34),   # KEPATUHAN - Orange
    RGBColor(52, 152, 219),   # PROSES - Light Blue
    RGBColor(46, 204, 113),   # DISETUJUI - Green
    RGBColor(231, 76, 60),    # DITOLAK - Red
]

margin = Emu(457200)  # 0.5 inch margin

# ============================================
# HELPER FUNCTIONS
# ============================================
def add_text(slide, text, left, top, width, height, size=Pt(14), bold=False, color=dark_blue, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    return box

def add_card(slide, title, desc, left, top, width, height, bg_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = bg_color
    tf = shape.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = white
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(9)
    p2.font.bold = False
    p2.font.color.rgb = RGBColor(220, 220, 220)
    return shape

def add_step_card(slide, num, title, desc, left, top, width, height):
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = light_bg
    bg.line.color.rgb = RGBColor(200, 210, 220)
    
    # Number circle
    circle_size = Emu(274320)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Emu(73152), top + Emu(73152), circle_size, circle_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = dark_blue
    circle.line.color.rgb = dark_blue
    tf_c = circle.text_frame
    tf_c.paragraphs[0].text = str(num)
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_c.paragraphs[0].font.size = Pt(12)
    tf_c.paragraphs[0].font.bold = True
    tf_c.paragraphs[0].font.color.rgb = white
    
    # Title
    add_text(slide, title,
             left + Emu(384048), top + Emu(73152),
             width - Emu(457200), Emu(228600),
             size=Pt(12), bold=True, color=dark_blue)
    
    # Description
    add_text(slide, desc,
             left + Emu(73152), top + Emu(329184),
             width - Emu(146304), height - Emu(365760),
             size=Pt(9), bold=False, color=grey_text)

# ============================================
# DARK SIDEBAR
# ============================================
sidebar = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(73152), Emu(SLIDE_H))
sidebar.fill.solid()
sidebar.fill.fore_color.rgb = dark_blue
sidebar.line.color.rgb = dark_blue

# ============================================
# HEADER: Number badge + Title
# ============================================
badge = slide5.shapes.add_shape(MSO_SHAPE.OVAL, margin, Emu(274320), Emu(548640), Emu(548640))
badge.fill.solid()
badge.fill.fore_color.rgb = gold
badge.line.color.rgb = gold
tf_b = badge.text_frame
tf_b.paragraphs[0].text = "2"
tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER
tf_b.paragraphs[0].font.size = Pt(22)
tf_b.paragraphs[0].font.bold = True
tf_b.paragraphs[0].font.color.rgb = white

add_text(slide5, "CARA KERJA ANALIS (INPUT KREDIT)",
         Emu(1143000), Emu(228600), Emu(9144000), Emu(411480),
         size=Pt(22), bold=True, color=dark_blue)

add_text(slide5, "Proses Input Data Pengajuan Kredit oleh Analis",
         Emu(1143000), Emu(640080), Emu(9144000), Emu(274320),
         size=Pt(13), bold=False, color=grey_text)

# Divider line
line = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, Emu(960120), Emu(SLIDE_W - 2 * 457200), Emu(27432))
line.fill.solid()
line.fill.fore_color.rgb = gold
line.line.color.rgb = gold

# ============================================
# SECTION 1: JENIS KREDIT
# ============================================
add_text(slide5, "JENIS KREDIT YANG TERSEDIA:",
         margin, Emu(1051560), Emu(10058400), Emu(274320),
         size=Pt(13), bold=True, color=dark_blue)

kredit_data = [
    ("Umum", "Kredit umum untuk\nnasabah perorangan"),
    ("PPPK", "Kredit untuk Pegawai\nPemerintah (PPPK)"),
    ("Perangkat Desa", "Kredit khusus\nPerangkat Desa"),
    ("KPR", "Kredit Pemilikan\nRumah"),
    ("Kretamas", "Kredit Emas\n(Gold Loan)"),
    ("Cash Collateral", "Kredit dengan\njaminan deposito"),
]

card_w = Emu(1783080)  # ~1.95 inches
card_h = Emu(640080)   # ~0.70 inches
card_y = Emu(1371600)
card_gap = Emu(91440)

for i, (title, desc) in enumerate(kredit_data):
    x = Emu(457200 + i * (card_w + card_gap))
    add_card(slide5, title, desc, x, card_y, card_w, card_h, card_colors[i])

# ============================================
# PENTING: Keterangan Plafond 50 Juta
# ============================================
note_y = Emu(2103120)
note_bg = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, margin, note_y, Emu(SLIDE_W - 2 * 457200), Emu(320040))
note_bg.fill.solid()
note_bg.fill.fore_color.rgb = RGBColor(253, 237, 236)
note_bg.line.color.rgb = RGBColor(231, 76, 60)

tf_note = note_bg.text_frame
tf_note.word_wrap = True
p_note = tf_note.paragraphs[0]
p_note.text = "\u26a0\ufe0f  PENTING: Jenis kredit di atas terbatas untuk maksimal plafond kredit Rp 50.000.000 (lima puluh juta rupiah)"
p_note.alignment = PP_ALIGN.CENTER
p_note.font.size = Pt(11)
p_note.font.bold = True
p_note.font.color.rgb = red_text

# ============================================
# SECTION 2: 6 TAHAP INPUT DATA
# ============================================
add_text(slide5, "6 TAHAP INPUT DATA KREDIT:",
         margin, Emu(2514600), Emu(10058400), Emu(274320),
         size=Pt(13), bold=True, color=dark_blue)

steps_data = [
    (1, "Data Pemohon", "NIK, Nama, Alamat, Pekerjaan, Tanggal Lahir, Status Perkawinan"),
    (2, "Penghasilan", "Gaji pokok, tunjangan, penghasilan lain, analisa repayment capacity"),
    (3, "Agunan / Jaminan", "Data kendaraan/tanah/bangunan, foto agunan, taksasi nilai"),
    (4, "Struktur Kredit", "Plafond, jangka waktu, suku bunga, tujuan penggunaan kredit"),
    (5, "Neraca Keuangan", "Aset, kewajiban, modal bersih (untuk kredit usaha/umum)"),
    (6, "Analisa 6C", "Character, Capacity, Capital, Collateral, Condition, Constraint"),
]

step_w = Emu(3566160)   # ~3.9 inches
step_h = Emu(640080)    # ~0.7 inches
step_gap_x = Emu(182880)
step_gap_y = Emu(137160)
step_start_y = Emu(2834640)

for i, (num, title, desc) in enumerate(steps_data):
    col = i % 3
    row = i // 3
    x = Emu(457200 + col * (step_w + step_gap_x))
    y = Emu(step_start_y + row * (step_h + step_gap_y))
    add_step_card(slide5, num, title, desc, x, y, step_w, step_h)

# ============================================
# SECTION 3: STATUS PENGAJUAN KREDIT
# ============================================
status_section_y = Emu(4480560)
add_text(slide5, "STATUS PENGAJUAN KREDIT (TERMASUK KEPATUHAN):",
         margin, status_section_y, Emu(10058400), Emu(274320),
         size=Pt(13), bold=True, color=dark_blue)

statuses = [
    ("DRAFT", "Analis buat\npengajuan baru"),
    ("DIAJUKAN", "Dikirim ke rantai\napproval"),
    ("KEPATUHAN", "Review kepatuhan\n& SOP"),
    ("PROSES", "Ditinjau pejabat\nberwenang"),
    ("DISETUJUI", "Kredit disetujui\noleh pejabat"),
    ("DITOLAK", "Kredit ditolak,\nkembali ke analis"),
]

stat_w = Emu(1783080)
stat_h = Emu(548640)
stat_gap = Emu(91440)
stat_y = Emu(4800600)

for i, (label, desc) in enumerate(statuses):
    x = Emu(457200 + i * (stat_w + stat_gap))
    add_card(slide5, label, desc, x, stat_y, stat_w, stat_h, status_colors[i])

# Add flow arrows between status boxes
arrow_y = Emu(stat_y + stat_h // 2 - 54864)
for i in range(5):
    x = Emu(457200 + (i + 1) * stat_w + i * stat_gap + stat_gap // 2)
    add_text(slide5, "\u2192", x, arrow_y, stat_gap, Emu(182880),
             size=Pt(16), bold=True, color=gold, align=PP_ALIGN.CENTER)

# ============================================
# SECTION 4: Alur singkat
# ============================================
alur_y = Emu(5486400)
alur_bg = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, margin, alur_y, Emu(SLIDE_W - 2 * 457200), Emu(411480))
alur_bg.fill.solid()
alur_bg.fill.fore_color.rgb = dark_blue
alur_bg.line.color.rgb = dark_blue

tf_alur = alur_bg.text_frame
tf_alur.word_wrap = True
p_alur = tf_alur.paragraphs[0]
p_alur.text = "ALUR APPROVAL:  ANALIS (Input)  \u2192  KASUBAG (Review)  \u2192  KEPATUHAN (Assessment)  \u2192  KABAG ANALIS (Evaluasi)  \u2192  KADIV (Approval)  \u2192  DIREKSI (Final)"
p_alur.alignment = PP_ALIGN.CENTER
p_alur.font.size = Pt(12)
p_alur.font.bold = True
p_alur.font.color.rgb = white

# ============================================
# SECTION 5: Keterangan tambahan
# ============================================
note2_y = Emu(5989320)
note2_bg = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, margin, note2_y, Emu(SLIDE_W - 2 * 457200), Emu(731520))
note2_bg.fill.solid()
note2_bg.fill.fore_color.rgb = light_bg
note2_bg.line.color.rgb = RGBColor(200, 210, 220)

tf_n2 = note2_bg.text_frame
tf_n2.word_wrap = True

info_items = [
    "\u2705 Kalkulasi Otomatis: Repayment capacity & taksasi agunan dihitung otomatis oleh sistem",
    "\u2705 Fitur Revisi: Pengajuan yang dikembalikan dapat diedit & dikirim ulang langsung ke level yang mengembalikan",
    "\u2705 Fitur Resubmit: Pengajuan ditolak dapat diperbaiki & diajukan kembali sebagai pengajuan baru",
]

p0 = tf_n2.paragraphs[0]
p0.text = info_items[0]
p0.font.size = Pt(10)
p0.font.color.rgb = dark_blue

for item in info_items[1:]:
    p = tf_n2.add_paragraph()
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = dark_blue

# ============================================
# VERIFY: Check nothing is offside
# ============================================
offside_count = 0
for shape in slide5.shapes:
    r = shape.left + shape.width
    b = shape.top + shape.height
    if r > SLIDE_W or b > SLIDE_H:
        offside_count += 1
        print(f'OFFSIDE: L={shape.left}, T={shape.top}, R={r}, B={b}')

if offside_count == 0:
    print("No offside shapes detected.")

prs.save(r'd:\laragon\www\andrian\Presentasi_Sistem_Kredit_BPR_Wonosobo.pptx')
print("Slide 5 berhasil didesain ulang secara presisi.")
