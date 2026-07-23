import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
prs = Presentation(r'd:\laragon\www\andrian\Presentasi_Sistem_Kredit_BPR_Wonosobo.pptx')

# Slide 5 (index 4)
slide5 = prs.slides[4]
with open('slide5_dump.txt', 'w', encoding='utf-8') as f:
    f.write('=== SLIDE 5 ===\n')
    for i, shape in enumerate(slide5.shapes):
        f.write(f'Shape {i}: type={shape.shape_type}, left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}\n')
        if shape.has_text_frame:
            for pi, p in enumerate(shape.text_frame.paragraphs):
                f.write(f'  P{pi}: "{p.text}"\n')
        f.write('\n')

# Slide 7 (index 6)
slide7 = prs.slides[6]
with open('slide7_dump.txt', 'w', encoding='utf-8') as f:
    f.write('=== SLIDE 7 ===\n')
    for i, shape in enumerate(slide7.shapes):
        f.write(f'Shape {i}: type={shape.shape_type}, left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}\n')
        if shape.has_text_frame:
            for pi, p in enumerate(shape.text_frame.paragraphs):
                f.write(f'  P{pi}: "{p.text}"\n')
        f.write('\n')

print('Done dumping.')
