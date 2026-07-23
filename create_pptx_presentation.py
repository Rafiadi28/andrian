#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Script untuk membuat PowerPoint Presentasi Sistem Aplikasi Kredit
PT BPR Bank Wonosobo
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# CONFIGURATION
# ============================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color palette
NAVY = RGBColor(0x0B, 0x1D, 0x3A)
DARK_BLUE = RGBColor(0x14, 0x2B, 0x52)
BLUE = RGBColor(0x1E, 0x5A, 0xA8)
LIGHT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
TEAL = RGBColor(0x06, 0xB6, 0xD4)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
DARK_GREEN = RGBColor(0x05, 0x9E, 0x6F)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY = RGBColor(0x47, 0x55, 0x69)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)

# Image paths
IMG_DIR = r"C:\Users\LAB RPL 18\.gemini\antigravity\brain\542e5a00-d231-49cf-b4fc-a2cdd75b32bc"
IMAGES = {
    'system_flow': os.path.join(IMG_DIR, 'system_flow_diagram_1784769099340.jpg'),
    'analis': os.path.join(IMG_DIR, 'analis_workflow_1784769106740.jpg'),
    'approval': os.path.join(IMG_DIR, 'approval_workflow_1784769114724.jpg'),
    'admin': os.path.join(IMG_DIR, 'admin_workflow_1784769134628.jpg'),
    'direksi': os.path.join(IMG_DIR, 'direksi_workflow_1784769143726.jpg'),
    'kadiv': os.path.join(IMG_DIR, 'kadiv_workflow_1784769151038.jpg'),
}

OUTPUT_PATH = r"d:\laragon\www\andrian\Presentasi_Sistem_Kredit_BPR_Wonosobo.pptx"

# ============================================================
# HELPER FUNCTIONS
# ============================================================

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a shape with optional fill and line."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14, font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2), font_name='Calibri'):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_card(slide, left, top, width, height, title, description, icon_text, card_color, accent_color):
    """Add a styled card with icon, title, and description."""
    # Card background
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color=RGBColor(0x1E, 0x29, 0x3B))
    card.line.color.rgb = RGBColor(0x33, 0x44, 0x55)
    card.line.width = Pt(1)

    # Accent bar on top
    add_shape(slide, MSO_SHAPE.RECTANGLE, left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), Inches(0.06), fill_color=accent_color)

    # Icon circle
    icon_size = Inches(0.55)
    icon = add_shape(slide, MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.25), icon_size, icon_size, fill_color=accent_color)
    icon_tf = icon.text_frame
    icon_tf.paragraphs[0].text = icon_text
    icon_tf.paragraphs[0].font.size = Pt(16)
    icon_tf.paragraphs[0].font.color.rgb = WHITE
    icon_tf.paragraphs[0].font.bold = True
    icon_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon_tf.word_wrap = False

    # Title
    add_textbox(slide, left + Inches(0.2), top + Inches(0.9), width - Inches(0.4), Inches(0.35),
                title, font_size=13, font_color=WHITE, bold=True)

    # Description
    add_textbox(slide, left + Inches(0.2), top + Inches(1.2), width - Inches(0.4), height - Inches(1.4),
                description, font_size=10, font_color=GRAY)

    return card


def add_step_box(slide, left, top, width, height, number, title, desc, color):
    """Add a numbered step box."""
    # Background
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color=RGBColor(0x1E, 0x29, 0x3B))
    box.line.color.rgb = color
    box.line.width = Pt(2)

    # Number badge
    badge_size = Inches(0.4)
    badge = add_shape(slide, MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.15), badge_size, badge_size, fill_color=color)
    badge_tf = badge.text_frame
    badge_tf.paragraphs[0].text = str(number)
    badge_tf.paragraphs[0].font.size = Pt(14)
    badge_tf.paragraphs[0].font.color.rgb = WHITE
    badge_tf.paragraphs[0].font.bold = True
    badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    add_textbox(slide, left + Inches(0.65), top + Inches(0.15), width - Inches(0.85), Inches(0.35),
                title, font_size=12, font_color=WHITE, bold=True)

    # Description
    add_textbox(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), height - Inches(0.65),
                desc, font_size=9, font_color=GRAY)


def add_decision_badge(slide, left, top, width, height, label, color, desc):
    """Add a decision badge (Setuju/Revisi/Tolak)."""
    badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color=color)
    tf = badge.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Description below
    add_textbox(slide, left, top + height + Inches(0.1), width, Inches(0.5),
                desc, font_size=9, font_color=GRAY, alignment=PP_ALIGN.CENTER)


def add_arrow_right(slide, left, top, width, color=GOLD):
    """Add a right arrow."""
    arrow = add_shape(slide, MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.3), fill_color=color)
    return arrow


def add_section_header(slide, number, title, subtitle, color):
    """Add a section header with number badge."""
    # Large number circle
    badge = add_shape(slide, MSO_SHAPE.OVAL, Inches(0.8), Inches(0.5), Inches(1.0), Inches(1.0), fill_color=color)
    tf = badge.text_frame
    tf.paragraphs[0].text = str(number)
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    add_textbox(slide, Inches(2.0), Inches(0.5), Inches(10.0), Inches(0.6),
                title, font_size=30, font_color=WHITE, bold=True)

    # Subtitle
    add_textbox(slide, Inches(2.0), Inches(1.1), Inches(10.0), Inches(0.4),
                subtitle, font_size=14, font_color=GOLD)

    # Divider line
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.04), fill_color=color)


# ============================================================
# SLIDE CREATION
# ============================================================

def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ========================================
    # SLIDE 1: TITLE SLIDE
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, NAVY)

    # Decorative accent shapes
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, fill_color=GOLD)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=GOLD)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_HEIGHT - Inches(0.06), SLIDE_WIDTH, Inches(0.06), fill_color=GOLD)

    # Title
    add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10.0), Inches(1.0),
                "SISTEM APLIKASI KREDIT", font_size=44, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.0), Inches(0.7),
                "PT BPR BANK WONOSOBO", font_size=28, font_color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    # Divider
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.4), Inches(4.3), Inches(0.04), fill_color=GOLD)

    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10.0), Inches(0.5),
                "Panduan Lengkap Alur Sistem & Cara Kerja Setiap Role", font_size=18, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Date
    add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10.0), Inches(0.5),
                "Presentasi Demo Sistem — Juli 2026", font_size=14, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # ========================================
    # SLIDE 2: AGENDA / DAFTAR ISI
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, fill_color=GOLD)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.0), Inches(0.7),
                "DAFTAR ISI PRESENTASI", font_size=32, font_color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(3.0), Inches(0.04), fill_color=GOLD)

    agenda_items = [
        ("1", "Alur Sistem", "Overview lengkap alur pengajuan hingga persetujuan", BLUE),
        ("2", "Cara Kerja Analis", "Proses input data kredit oleh analis", TEAL),
        ("3", "Cara Kerja Kepatuhan", "Penilaian kepatuhan & kelengkapan dokumen", DARK_GREEN),
        ("4", "Cara Kerja Kasubag", "Review & verifikasi awal oleh Kasubag", GREEN),
        ("5", "Cara Kerja Kabag", "Evaluasi risiko kredit oleh Kabag", ORANGE),
        ("6", "Cara Kerja Kadiv", "Approval divisi & threshold keputusan", PURPLE),
        ("7", "Cara Kerja Direksi", "Persetujuan final untuk kredit besar", RED),
        ("8", "Cara Kerja Admin", "Manajemen sistem, user & parameter", GOLD),
    ]

    for i, (num, title, desc, color) in enumerate(agenda_items):
        y = Inches(1.4) + Inches(i * 0.7)
        # Number badge
        badge = add_shape(slide, MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.5), Inches(0.5), fill_color=color)
        tf = badge.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, Inches(1.7), y + Inches(0.02), Inches(6.0), Inches(0.3),
                    title, font_size=16, font_color=WHITE, bold=True)
        add_textbox(slide, Inches(1.7), y + Inches(0.32), Inches(8.0), Inches(0.3),
                    desc, font_size=11, font_color=GRAY)

    # ========================================
    # SLIDE 3: ALUR SISTEM - OVERVIEW
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "1", "ALUR SISTEM APLIKASI KREDIT", "Overview Proses Pengajuan Hingga Persetujuan", BLUE)

    # Flow diagram using shapes
    roles_flow = [
        ("ANALIS\n(Input)", BLUE),
        ("KEPATUHAN\n(Review)", DARK_GREEN),
        ("KASUBAG\n(Review)", TEAL),
        ("KABAG\n(Evaluasi)", ORANGE),
        ("KADIV\n(Approval)", PURPLE),
        ("DIREKSI\n(Final)", RED),
    ]

    box_w = Inches(1.6)
    box_h = Inches(1.0)
    start_x = Inches(0.5)
    y_pos = Inches(2.2)
    gap = Inches(0.4)

    for i, (label, color) in enumerate(roles_flow):
        x = start_x + i * (box_w + gap)
        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, box_w, box_h, fill_color=color)
        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        if i < len(roles_flow) - 1:
            arrow_x = x + box_w + Inches(0.05)
            add_arrow_right(slide, arrow_x, y_pos + Inches(0.35), Inches(0.45), GOLD)

    # Status descriptions below
    statuses = [
        ("DRAFT", "Analis membuat pengajuan kredit baru", BLUE),
        ("DIAJUKAN", "Pengajuan dikirim ke rantai approval", TEAL),
        ("PROSES", "Sedang ditinjau oleh pejabat berwenang", ORANGE),
        ("DISETUJUI", "Kredit disetujui oleh pejabat terakhir", GREEN),
        ("DITOLAK", "Kredit ditolak, dikembalikan ke analis", RED),
        ("REVISI", "Perlu perbaikan data oleh analis", GOLD),
    ]

    add_textbox(slide, Inches(0.8), Inches(3.6), Inches(11.0), Inches(0.4),
                "STATUS PENGAJUAN KREDIT:", font_size=14, font_color=GOLD, bold=True)

    for i, (status, desc, color) in enumerate(statuses):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + col * Inches(4.0)
        y = Inches(4.1) + row * Inches(0.7)

        badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.2), Inches(0.4), fill_color=color)
        tf = badge.text_frame
        tf.paragraphs[0].text = status
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, x + Inches(1.35), y + Inches(0.05), Inches(2.5), Inches(0.35),
                    desc, font_size=10, font_color=LIGHT_GRAY)

    # Auto-skip note
    note_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.2), fill_color=DARK_BLUE)
    note_box.line.color.rgb = GOLD
    note_box.line.width = Pt(1)

    txBox = add_textbox(slide, Inches(1.0), Inches(5.7), Inches(11.0), Inches(0.3),
                "⚡ FITUR AUTO-SKIP & THRESHOLD AMOUNT", font_size=13, font_color=GOLD, bold=True)

    txBox2 = add_textbox(slide, Inches(1.0), Inches(6.05), Inches(11.0), Inches(0.7),
                "• Jika pejabat sedang cuti/sakit/berhalangan, sistem otomatis melompati ke pejabat aktif berikutnya\n"
                "• Kredit < Rp 500 Juta: Approval final di level KADIV (tidak perlu ke Direksi)\n"
                "• Kredit ≥ Rp 500 Juta: Harus mendapat persetujuan DIREKSI sebagai level tertinggi",
                font_size=11, font_color=LIGHT_GRAY)

    # ========================================
    # SLIDE 4: ALUR SISTEM - DIAGRAM IMAGE
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, fill_color=BLUE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.0), Inches(0.5),
                "DIAGRAM ALUR APPROVAL KREDIT", font_size=28, font_color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.9), Inches(3.0), Inches(0.04), fill_color=BLUE)

    # Add system flow image
    if os.path.exists(IMAGES['system_flow']):
        slide.shapes.add_picture(IMAGES['system_flow'], Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))

    # ========================================
    # SLIDE 5: CARA KERJA ANALIS - OVERVIEW
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "2", "CARA KERJA ANALIS (INPUT KREDIT)", "Proses Input Data Pengajuan Kredit oleh Analis", TEAL)

    # Credit types
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.0), Inches(0.4),
                "JENIS KREDIT YANG TERSEDIA:", font_size=14, font_color=GOLD, bold=True)

    credit_types = [
        ("Umum", "Kredit umum untuk\nnasabah perorangan", BLUE),
        ("PPPK", "Kredit untuk Pegawai\nPemerintah (PPPK)", TEAL),
        ("Perangkat Desa", "Kredit khusus\nPerangkat Desa", GREEN),
        ("KPR", "Kredit Pemilikan\nRumah", ORANGE),
        ("Kretamas", "Kredit Emas\n(Gold Loan)", GOLD),
        ("Cash Collateral", "Kredit dengan\njaminan deposito", PURPLE),
    ]

    for i, (name, desc, color) in enumerate(credit_types):
        x = Inches(0.8) + (i % 6) * Inches(2.05)
        y = Inches(2.5)
        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.9), Inches(1.1), fill_color=DARK_BLUE)
        box.line.color.rgb = color
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = name
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # 6 Input Sections
    add_textbox(slide, Inches(0.8), Inches(3.9), Inches(11.0), Inches(0.4),
                "6 TAHAP INPUT DATA KREDIT:", font_size=14, font_color=GOLD, bold=True)

    steps = [
        ("1", "Data Pemohon", "NIK, Nama, Alamat, Pekerjaan,\nTanggal Lahir, Status Perkawinan", BLUE),
        ("2", "Penghasilan", "Gaji pokok, tunjangan, penghasilan\nlain, analisa repayment capacity", TEAL),
        ("3", "Agunan / Jaminan", "Data kendaraan/tanah/bangunan,\nfoto agunan, taksasi nilai", GREEN),
        ("4", "Struktur Kredit", "Plafond, jangka waktu, suku bunga,\ntujuan penggunaan kredit", ORANGE),
        ("5", "Neraca Keuangan", "Aset, kewajiban, modal bersih\n(untuk kredit usaha/umum)", PURPLE),
        ("6", "Analisa 6C & Verifikasi", "Character, Capacity, Capital,\nCollateral, Condition, Constraint", RED),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + col * Inches(4.1)
        y = Inches(4.4) + row * Inches(1.35)
        add_step_box(slide, x, y, Inches(3.8), Inches(1.2), num, title, desc, color)

    # ========================================
    # SLIDE 6: ANALIS - WORKFLOW IMAGE
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, fill_color=TEAL)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.0), Inches(0.5),
                "WORKFLOW INPUT DATA ANALIS", font_size=28, font_color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.9), Inches(3.0), Inches(0.04), fill_color=TEAL)

    if os.path.exists(IMAGES['analis']):
        slide.shapes.add_picture(IMAGES['analis'], Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))

    # ========================================
    # SLIDE 7: ANALIS - DETAIL PROSES
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, fill_color=TEAL)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.0), Inches(0.5),
                "DETAIL PROSES KERJA ANALIS", font_size=28, font_color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.9), Inches(3.0), Inches(0.04), fill_color=TEAL)

    # Left column - Flow steps
    flow_steps = [
        ("Login ke Sistem", "Analis login menggunakan akun yang terdaftar\ndengan role 'analis'"),
        ("Pilih Jenis Kredit", "Memilih jenis produk kredit sesuai kebutuhan\nnasabah (Umum/PPPK/Desa/dll)"),
        ("Input Data Lengkap", "Mengisi form data pemohon, penghasilan,\nagunan, struktur kredit, neraca, & analisa 6C"),
        ("Upload Dokumen", "Upload foto agunan, dokumen pendukung,\ndan kelengkapan administrasi"),
        ("Simpan Draft", "Data dapat disimpan sebagai draft untuk\ndilengkapi kemudian"),
        ("Submit Pengajuan", "Kirim pengajuan ke rantai approval.\nStatus berubah menjadi 'Diajukan'"),
    ]

    for i, (title, desc) in enumerate(flow_steps):
        y = Inches(1.3) + i * Inches(0.95)
        add_step_box(slide, Inches(0.5), y, Inches(5.8), Inches(0.85), i + 1, title, desc, TEAL)

    # Right column - Special features
    right_x = Inches(6.8)
    add_textbox(slide, right_x, Inches(1.3), Inches(6.0), Inches(0.4),
                "FITUR KHUSUS ANALIS:", font_size=14, font_color=GOLD, bold=True)

    features = [
        ("📝 Revisi", "Jika pengajuan dikembalikan untuk revisi,\nanalis dapat mengedit data dan mengirim ulang.\nSistem langsung mengirim ke level yang\nmengembalikan (skip level sebelumnya).", ORANGE),
        ("📊 Kalkulasi Otomatis", "Repayment capacity, taksasi agunan,\ndan rasio plafond dihitung otomatis\nberdasarkan parameter yang dikonfigurasi.", BLUE),
        ("📋 Riwayat", "Melihat seluruh riwayat pengajuan yang\npernah dibuat beserta status terkini.", PURPLE),
        ("🔄 Resubmit", "Pengajuan yang ditolak dapat diperbaiki\ndan diajukan kembali sebagai pengajuan baru.", RED),
    ]

    for i, (title, desc, color) in enumerate(features):
        y = Inches(1.8) + i * Inches(1.3)
        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, right_x, y, Inches(5.8), Inches(1.15), fill_color=DARK_BLUE)
        box.line.color.rgb = color
        box.line.width = Pt(1.5)

        add_textbox(slide, right_x + Inches(0.2), y + Inches(0.1), Inches(5.4), Inches(0.3),
                    title, font_size=13, font_color=color, bold=True)
        add_textbox(slide, right_x + Inches(0.2), y + Inches(0.4), Inches(5.4), Inches(0.7),
                    desc, font_size=10, font_color=LIGHT_GRAY)

    # ========================================
    # SLIDE 8: CARA KERJA KEPATUHAN
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "3", "CARA KERJA KEPATUHAN", "Singkronisasi & Penilaian Dokumen Kepatuhan", DARK_GREEN)

    # Workflow
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.0), Inches(0.4),
                "ALUR KERJA KEPATUHAN:", font_size=14, font_color=GOLD, bold=True)

    kepatuhan_steps = [
        ("1", "Terima Data Analis", "Sistem auto-populate form dari\ndata input awal Analis", BLUE),
        ("2", "Review Kepatuhan", "Periksa dokumen, checklist SOP,\ndan validasi syarat kredit", DARK_GREEN),
        ("3", "Singkronisasi Data", "Melengkapi fasilitas kredit,\ncatatan, & kesimpulan final", TEAL),
        ("4", "Simpan Assessment", "Data tersimpan di database,\nsiap direview oleh Komite", ORANGE),
    ]

    for i, (num, title, desc, color) in enumerate(kepatuhan_steps):
        x = Inches(0.8) + i * Inches(3.1)
        add_step_box(slide, x, Inches(2.5), Inches(2.9), Inches(1.2), num, title, desc, color)

        if i < len(kepatuhan_steps) - 1:
            add_arrow_right(slide, x + Inches(2.95), Inches(2.95), Inches(0.1), GOLD)

    # Focus areas
    add_textbox(slide, Inches(0.8), Inches(4.1), Inches(11.0), Inches(0.4),
                "FITUR UTAMA PENILAIAN KEPATUHAN:", font_size=14, font_color=GOLD, bold=True)

    focus_items_kepatuhan = [
        ("📝", "Auto-Populate Data", "Data dari Analis langsung muncul\ndi form Kepatuhan (tidak perlu\ninput ulang).", TEAL),
        ("🔄", "AJAX Submission", "Penyimpanan data real-time\ntanpa loading ulang halaman.", DARK_GREEN),
        ("🔍", "Validasi Dokumen", "Memastikan semua syarat legal,\nagunan, dan administrasi telah\nsesuai SOP Bank.", ORANGE),
        ("📊", "Tersentralisasi", "Hasil assessment langsung\nterhubung ke pengajuan kredit\ndan bisa dilihat Kasubag/Kabag.", BLUE),
    ]

    for i, (icon, title, desc, color) in enumerate(focus_items_kepatuhan):
        x = Inches(0.8) + i * Inches(3.1)
        y = Inches(4.6)

        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(2.3), fill_color=DARK_BLUE)
        box.line.color.rgb = color
        box.line.width = Pt(1.5)

        # Icon
        add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.4),
                    icon, font_size=22, font_color=WHITE, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.5), Inches(0.3),
                    title, font_size=13, font_color=color, bold=True)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.5), Inches(1.2),
                    desc, font_size=10, font_color=LIGHT_GRAY)

    # ========================================
    # SLIDE 9: CARA KERJA KASUBAG
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "4", "CARA KERJA KASUBAG ANALIS", "Review & Verifikasi Awal Pengajuan Kredit", GREEN)

    # Main workflow
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.0), Inches(0.4),
                "ALUR KERJA KASUBAG ANALIS:", font_size=14, font_color=GOLD, bold=True)

    kasubag_steps = [
        ("1", "Buka Dashboard", "Login dan lihat daftar pengajuan\nyang menunggu review di inbox", GREEN),
        ("2", "Review Detail", "Klik pengajuan untuk melihat\nseluruh data yang diinput analis", TEAL),
        ("3", "Verifikasi Dokumen", "Periksa kelengkapan data pemohon,\nkesesuaian dokumen pendukung", BLUE),
        ("4", "Cek Perhitungan", "Verifikasi kalkulasi repayment\ncapacity dan taksasi agunan", PURPLE),
    ]

    for i, (num, title, desc, color) in enumerate(kasubag_steps):
        x = Inches(0.8) + i * Inches(3.1)
        add_step_box(slide, x, Inches(2.5), Inches(2.9), Inches(1.2), num, title, desc, color)

        if i < len(kasubag_steps) - 1:
            add_arrow_right(slide, x + Inches(2.95), Inches(2.95), Inches(0.1), GOLD)

    # Decision section
    add_textbox(slide, Inches(0.8), Inches(4.1), Inches(11.0), Inches(0.4),
                "KEPUTUSAN YANG DAPAT DIAMBIL:", font_size=14, font_color=GOLD, bold=True)

    decisions = [
        ("✅ SETUJU", GREEN, "Pengajuan diteruskan ke\nKABAG KREDIT\nuntuk evaluasi lebih lanjut",
         "Status: posisi_saat_ini → 'kabag_kredit'"),
        ("🔄 REVISI", ORANGE, "Pengajuan dikembalikan ke\nANALIS\nuntuk perbaikan data/dokumen",
         "Status: posisi_saat_ini → 'analis' | status → 'revisi'"),
        ("❌ TOLAK", RED, "Pengajuan DITOLAK dan\ndikembalikan ke ANALIS\ndengan alasan penolakan",
         "Status: posisi_saat_ini → 'analis' | status → 'ditolak'"),
    ]

    for i, (label, color, desc, status) in enumerate(decisions):
        x = Inches(0.8) + i * Inches(4.1)
        y = Inches(4.6)

        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.3), fill_color=DARK_BLUE)
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Decision badge
        badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.8), y + Inches(0.2), Inches(2.2), Inches(0.5), fill_color=color)
        tf = badge.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, x + Inches(0.2), y + Inches(0.85), Inches(3.4), Inches(0.8),
                    desc, font_size=11, font_color=WHITE, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.2), y + Inches(1.7), Inches(3.4), Inches(0.45),
                    status, font_size=8, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    # ========================================
    # SLIDE 9: CARA KERJA KABAG
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "5", "CARA KERJA KABAG KREDIT", "Evaluasi Risiko Kredit & Penilaian Agunan", ORANGE)

    # Workflow
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.0), Inches(0.4),
                "ALUR KERJA KABAG KREDIT:", font_size=14, font_color=GOLD, bold=True)

    kabag_steps = [
        ("1", "Dashboard Inbox", "Menerima pengajuan yang sudah\ndisetujui oleh Kasubag Analis", ORANGE),
        ("2", "Evaluasi Risiko", "Analisis mendalam terhadap risiko\nkredit, kelayakan usaha nasabah", RED),
        ("3", "Review Agunan", "Verifikasi nilai taksasi agunan\ndan kecukupan jaminan", PURPLE),
        ("4", "Keputusan", "Memberikan keputusan: Setuju,\nRevisi, atau Tolak", GREEN),
    ]

    for i, (num, title, desc, color) in enumerate(kabag_steps):
        x = Inches(0.8) + i * Inches(3.1)
        add_step_box(slide, x, Inches(2.5), Inches(2.9), Inches(1.2), num, title, desc, color)

        if i < len(kabag_steps) - 1:
            add_arrow_right(slide, x + Inches(2.95), Inches(2.95), Inches(0.1), GOLD)

    # Focus areas
    add_textbox(slide, Inches(0.8), Inches(4.1), Inches(11.0), Inches(0.4),
                "FOKUS EVALUASI KABAG KREDIT:", font_size=14, font_color=GOLD, bold=True)

    focus_items = [
        ("💰", "Analisis Keuangan", "Mengevaluasi rasio keuangan nasabah,\nDebt Service Ratio (DSR), dan\nRepayment Capacity sesuai standar bank", BLUE),
        ("🏠", "Penilaian Agunan", "Memverifikasi nilai taksasi agunan,\nrasio Loan to Value (LTV),\ndan kelayakan jaminan yang diajukan", GREEN),
        ("📊", "Risiko Kredit", "Menilai risiko kredit berdasarkan\nanalisa 6C yang dilakukan analis,\nkelayakan struktur kredit", ORANGE),
        ("📋", "Kepatuhan", "Memastikan kesesuaian dengan\nkebijakan internal bank dan\nregulasi OJK yang berlaku", PURPLE),
    ]

    for i, (icon, title, desc, color) in enumerate(focus_items):
        x = Inches(0.8) + i * Inches(3.1)
        y = Inches(4.6)

        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(2.3), fill_color=DARK_BLUE)
        box.line.color.rgb = color
        box.line.width = Pt(1.5)

        # Icon
        add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.4),
                    icon, font_size=22, font_color=WHITE, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.5), Inches(0.3),
                    title, font_size=13, font_color=color, bold=True)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.5), Inches(1.2),
                    desc, font_size=10, font_color=LIGHT_GRAY)

    # ========================================
    # SLIDE 10: KABAG - KEPUTUSAN (with approval image)
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, fill_color=ORANGE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.0), Inches(0.5),
                "PROSES KEPUTUSAN APPROVAL (Kasubag / Kabag)", font_size=26, font_color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.9), Inches(3.0), Inches(0.04), fill_color=ORANGE)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.0), Inches(0.5),
                "Tampilan halaman review dan aksi keputusan yang tersedia pada setiap level approval:", font_size=12, font_color=GRAY)

    if os.path.exists(IMAGES['approval']):
        slide.shapes.add_picture(IMAGES['approval'], Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.5))

    # ========================================
    # SLIDE 11: CARA KERJA KADIV
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "6", "CARA KERJA KADIV BISNIS", "Approval Divisi & Threshold Keputusan Kredit", PURPLE)

    # Key highlight - Amount threshold
    highlight = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1), Inches(11.7), Inches(1.5), fill_color=DARK_BLUE)
    highlight.line.color.rgb = GOLD
    highlight.line.width = Pt(2)

    add_textbox(slide, Inches(1.0), Inches(2.2), Inches(11.0), Inches(0.4),
                "⚡ ATURAN THRESHOLD JUMLAH KREDIT:", font_size=16, font_color=GOLD, bold=True)

    # Left - < 500 Juta
    box1 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(2.7), Inches(5.0), Inches(0.7), fill_color=GREEN)
    tf = box1.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Kredit < Rp 500 Juta → KADIV = FINAL APPROVAL"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Right - >= 500 Juta
    box2 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2.7), Inches(5.5), Inches(0.7), fill_color=RED)
    tf = box2.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "Kredit ≥ Rp 500 Juta → Eskalasi ke DIREKSI"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Workflow steps
    add_textbox(slide, Inches(0.8), Inches(3.9), Inches(11.0), Inches(0.4),
                "ALUR KERJA KADIV BISNIS:", font_size=14, font_color=GOLD, bold=True)

    kadiv_steps = [
        ("1", "Terima Pengajuan", "Menerima pengajuan dari Kabag\nKredit yang sudah disetujui", PURPLE),
        ("2", "Review Menyeluruh", "Evaluasi kelayakan bisnis, risiko\ndivisi, dan dampak portfolio", BLUE),
        ("3", "Cek Threshold", "Periksa jumlah kredit untuk\nmenentukan apakah Kadiv\nadalah level approval terakhir", GOLD),
        ("4", "Keputusan", "Setuju (final jika <500Jt),\nRevisi, atau Tolak.\nJika ≥500Jt: eskalasi ke Direksi", GREEN),
    ]

    for i, (num, title, desc, color) in enumerate(kadiv_steps):
        x = Inches(0.8) + i * Inches(3.1)
        add_step_box(slide, x, Inches(4.4), Inches(2.9), Inches(1.4), num, title, desc, color)

        if i < len(kadiv_steps) - 1:
            add_arrow_right(slide, x + Inches(2.95), Inches(4.9), Inches(0.1), GOLD)

    # Decision options
    add_textbox(slide, Inches(0.8), Inches(6.1), Inches(11.0), Inches(0.35),
                "OPSI KEPUTUSAN:  ✅ Setuju (lanjut/final)  |  🔄 Revisi (kembali ke Analis)  |  ❌ Tolak (ditolak ke Analis)",
                font_size=12, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ========================================
    # SLIDE 12: KADIV - DIAGRAM IMAGE
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, fill_color=PURPLE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.0), Inches(0.5),
                "DIAGRAM KEPUTUSAN KADIV BISNIS", font_size=28, font_color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.9), Inches(3.0), Inches(0.04), fill_color=PURPLE)

    if os.path.exists(IMAGES['kadiv']):
        slide.shapes.add_picture(IMAGES['kadiv'], Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))

    # ========================================
    # SLIDE 13: CARA KERJA DIREKSI
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "7", "CARA KERJA DIREKSI", "Persetujuan Final untuk Kredit Bernilai Besar (≥ Rp 500 Juta)", RED)

    # Important notice
    notice = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1), Inches(11.7), Inches(1.0), fill_color=RGBColor(0x45, 0x12, 0x12))
    notice.line.color.rgb = RED
    notice.line.width = Pt(2)

    add_textbox(slide, Inches(1.0), Inches(2.2), Inches(11.0), Inches(0.35),
                "⚠️  PENTING: Direksi hanya menerima pengajuan kredit dengan nilai ≥ Rp 500.000.000", font_size=14, font_color=RED, bold=True)
    add_textbox(slide, Inches(1.0), Inches(2.6), Inches(11.0), Inches(0.35),
                "Direksi TIDAK memiliki opsi REVISI — hanya dapat MENYETUJUI atau MENOLAK (keputusan bersifat FINAL)", font_size=12, font_color=RGBColor(0xFC, 0xA5, 0xA5))

    # Workflow
    add_textbox(slide, Inches(0.8), Inches(3.4), Inches(11.0), Inches(0.4),
                "ALUR KERJA DIREKSI:", font_size=14, font_color=GOLD, bold=True)

    direksi_steps = [
        ("1", "Dashboard Direksi", "Melihat pengajuan kredit bernilai\nbesar (≥500Jt) yang menunggu\nkeputusan eksekutif", RED),
        ("2", "Review Executive", "Evaluasi final terhadap kelayakan\nbisnis, risiko bank, dan dampak\nterhadap portfolio keseluruhan", BLUE),
        ("3", "Keputusan FINAL", "Hanya 2 pilihan:\n✅ SETUJU → Kredit Disetujui\n❌ TOLAK → Kredit Ditolak", GREEN),
    ]

    for i, (num, title, desc, color) in enumerate(direksi_steps):
        x = Inches(0.8) + i * Inches(4.1)
        add_step_box(slide, x, Inches(3.9), Inches(3.8), Inches(1.5), num, title, desc, color)

        if i < len(direksi_steps) - 1:
            add_arrow_right(slide, x + Inches(3.85), Inches(4.5), Inches(0.2), GOLD)

    # Two decision boxes
    add_textbox(slide, Inches(0.8), Inches(5.7), Inches(11.0), Inches(0.4),
                "HASIL KEPUTUSAN DIREKSI:", font_size=14, font_color=GOLD, bold=True)

    # SETUJU
    box_setuju = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.2), Inches(5.5), Inches(0.9), fill_color=GREEN)
    tf = box_setuju.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "✅ SETUJU — APPROVAL FINAL"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "status_pengajuan → 'disetujui' | posisi → 'selesai'"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0xBB, 0xF7, 0xD0)
    p.alignment = PP_ALIGN.CENTER

    # TOLAK
    box_tolak = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(6.2), Inches(5.5), Inches(0.9), fill_color=RED)
    tf = box_tolak.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "❌ TOLAK — PENOLAKAN FINAL"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "status_pengajuan → 'ditolak' | posisi → 'analis'"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0xFE, 0xCA, 0xCA)
    p.alignment = PP_ALIGN.CENTER

    # ========================================
    # SLIDE 14: DIREKSI - DIAGRAM IMAGE
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, fill_color=RED)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.0), Inches(0.5),
                "DIAGRAM WORKFLOW DIREKSI", font_size=28, font_color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.9), Inches(3.0), Inches(0.04), fill_color=RED)

    if os.path.exists(IMAGES['direksi']):
        slide.shapes.add_picture(IMAGES['direksi'], Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))

    # ========================================
    # SLIDE 15: CARA KERJA ADMIN
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_section_header(slide, "8", "CARA KERJA ADMIN / SUPERADMIN", "Manajemen Sistem, User, Parameter & Monitoring", GOLD)

    admin_features = [
        ("👥", "Manajemen User", "Membuat, mengedit, dan mengatur akun pengguna.\nMengatur role (Analis, Kasubag, Kabag, Kadiv, Direksi).\nMengatur status jabatan (Aktif, Cuti, Sakit, Izin, Berhalangan).", BLUE),
        ("🏛️", "Master Pejabat", "Konfigurasi jabatan dan tanda tangan pejabat\nyang tampil pada dokumen cetak.\nMenentukan urutan hierarki approval.", TEAL),
        ("⚙️", "Parameter Repayment", "Konfigurasi formula Repayment Capacity.\nAtur parameter DSR (Debt Service Ratio),\nbiaya hidup, dan threshold perhitungan.", GREEN),
        ("💾", "Backup Database", "Backup dan restore database sistem.\nMenjaga keamanan data dan\nmemastikan ketersediaan data.", ORANGE),
        ("🔍", "Audit Log", "Memonitor log aktivitas sistem.\nMelacak perubahan parameter.\nAudit trail untuk kepatuhan.", PURPLE),
        ("📊", "Riwayat & Monitoring", "Melihat seluruh riwayat pengajuan kredit.\nMonitor status pengajuan secara real-time.\nLaporan statistik dan aktivitas.", RED),
    ]

    for i, (icon, title, desc, color) in enumerate(admin_features):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + col * Inches(4.1)
        y = Inches(2.1) + row * Inches(2.5)

        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.3), fill_color=DARK_BLUE)
        box.line.color.rgb = color
        box.line.width = Pt(1.5)

        # Icon
        add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.45),
                    icon, font_size=24, font_color=WHITE)

        # Title
        add_textbox(slide, x + Inches(0.75), y + Inches(0.15), Inches(2.8), Inches(0.35),
                    title, font_size=15, font_color=color, bold=True)

        # Accent line
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.55), Inches(3.4), Inches(0.03), fill_color=color)

        # Description
        add_textbox(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.4), Inches(1.5),
                    desc, font_size=10, font_color=LIGHT_GRAY)

    # ========================================
    # SLIDE 16: ADMIN - DIAGRAM IMAGE
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, fill_color=GOLD)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.0), Inches(0.5),
                "FITUR ADMIN / SUPERADMIN", font_size=28, font_color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.9), Inches(3.0), Inches(0.04), fill_color=GOLD)

    if os.path.exists(IMAGES['admin']):
        slide.shapes.add_picture(IMAGES['admin'], Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8))

    # ========================================
    # SLIDE 17: ADMIN - DETAIL STATUS JABATAN
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, fill_color=GOLD)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.0), Inches(0.5),
                "MANAJEMEN STATUS JABATAN & AUTO-SKIP", font_size=26, font_color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.9), Inches(3.0), Inches(0.04), fill_color=GOLD)

    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.0), Inches(0.6),
                "Admin mengatur status jabatan setiap pejabat. Sistem secara otomatis melompati pejabat yang tidak aktif\ndan mengirim pengajuan ke pejabat aktif berikutnya dalam rantai approval.",
                font_size=12, font_color=LIGHT_GRAY)

    statuses = [
        ("AKTIF", "Pejabat aktif bertugas dan\ndapat menerima pengajuan kredit", GREEN),
        ("CUTI", "Pejabat sedang cuti,\npengajuan di-skip otomatis", ORANGE),
        ("SAKIT", "Pejabat sedang sakit,\npengajuan di-skip otomatis", RED),
        ("IZIN", "Pejabat sedang izin,\npengajuan di-skip otomatis", BLUE),
        ("BERHALANGAN", "Pejabat berhalangan,\npengajuan di-skip otomatis", PURPLE),
    ]

    for i, (status, desc, color) in enumerate(statuses):
        x = Inches(0.8) + i * Inches(2.45)
        y = Inches(2.1)

        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.3), Inches(1.6), fill_color=DARK_BLUE)
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Status badge
        badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), y + Inches(0.2), Inches(1.7), Inches(0.45), fill_color=color)
        tf = badge.text_frame
        tf.paragraphs[0].text = status
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, x + Inches(0.15), y + Inches(0.8), Inches(2.0), Inches(0.7),
                    desc, font_size=10, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Auto-skip diagram
    add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.0), Inches(0.4),
                "CONTOH SKENARIO AUTO-SKIP:", font_size=14, font_color=GOLD, bold=True)

    # Scenario flow
    scenario_roles = [
        ("ANALIS\n✅ Submit", BLUE, True),
        ("KASUBAG\n❌ CUTI", RED, False),
        ("KABAG\n✅ AKTIF", GREEN, True),
        ("KADIV\n❌ SAKIT", RED, False),
        ("DIREKSI\n✅ AKTIF", GREEN, True),
    ]

    for i, (label, color, active) in enumerate(scenario_roles):
        x = Inches(0.8) + i * Inches(2.5)
        y = Inches(4.5)

        bg_color = DARK_GREEN if active else RGBColor(0x45, 0x12, 0x12)
        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(1.0), fill_color=bg_color)
        box.line.color.rgb = color
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        if i < len(scenario_roles) - 1:
            arr_label = "SKIP →" if not active else "→"
            add_textbox(slide, x + Inches(2.2), y + Inches(0.3), Inches(0.3), Inches(0.4),
                        arr_label, font_size=9, font_color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.2),
                "Pada contoh di atas:\n"
                "• Analis submit pengajuan → Kasubag sedang CUTI → Sistem otomatis SKIP ke Kabag\n"
                "• Kabag AKTIF → Review & Setuju → Kadiv sedang SAKIT → Sistem otomatis SKIP ke Direksi\n"
                "• Direksi AKTIF → Review & memberikan keputusan FINAL\n"
                "• Setiap auto-skip dicatat dalam audit trail untuk transparansi",
                font_size=11, font_color=LIGHT_GRAY)

    # ========================================
    # SLIDE 18: RINGKASAN PERAN
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_HEIGHT, fill_color=GOLD)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.0), Inches(0.6),
                "RINGKASAN PERAN & TANGGUNG JAWAB", font_size=30, font_color=WHITE, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.9), Inches(4.0), Inches(0.04), fill_color=GOLD)

    summary_roles = [
        ("ANALIS", "Input data kredit, upload dokumen,\nsubmit & resubmit pengajuan", "Input & Submit", BLUE),
        ("KASUBAG", "Verifikasi awal, cek kelengkapan,\nreview data analis", "Review Awal", TEAL),
        ("KABAG", "Evaluasi risiko kredit, penilaian\nagunan, analisis keuangan", "Evaluasi Risiko", ORANGE),
        ("KADIV", "Approval divisi, final untuk <500Jt,\neskalasi untuk ≥500Jt", "Approval Divisi", PURPLE),
        ("DIREKSI", "Keputusan final untuk kredit ≥500Jt,\nhanya Setuju atau Tolak", "Final Approval", RED),
        ("ADMIN", "Manajemen user, parameter,\nbackup, audit & monitoring", "Sistem Admin", GOLD),
    ]

    for i, (role, desc, action, color) in enumerate(summary_roles):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + col * Inches(4.1)
        y = Inches(1.3) + row * Inches(2.9)

        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.6), fill_color=DARK_BLUE)
        box.line.color.rgb = color
        box.line.width = Pt(2)

        # Role name
        role_badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.5), fill_color=color)
        tf = role_badge.text_frame
        tf.paragraphs[0].text = role
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Action type
        add_textbox(slide, x + Inches(0.2), y + Inches(0.85), Inches(3.4), Inches(0.3),
                    action, font_size=13, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, x + Inches(0.2), y + Inches(1.2), Inches(3.4), Inches(1.2),
                    desc, font_size=11, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ========================================
    # SLIDE 19: PENUTUP
    # ========================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Decorative
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT, fill_color=GOLD)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=GOLD)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_HEIGHT - Inches(0.06), SLIDE_WIDTH, Inches(0.06), fill_color=GOLD)

    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.0), Inches(1.0),
                "TERIMA KASIH", font_size=48, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.2), Inches(4.3), Inches(0.04), fill_color=GOLD)

    add_textbox(slide, Inches(1.5), Inches(3.6), Inches(10.0), Inches(0.6),
                "Sistem Aplikasi Kredit — PT BPR Bank Wonosobo", font_size=20, font_color=GOLD, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10.0), Inches(0.5),
                "Presentasi Demo Sistem untuk Manajemen", font_size=16, font_color=GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10.0), Inches(0.5),
                "Siap untuk sesi tanya jawab & demo live", font_size=14, font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ========================================
    # SAVE
    # ========================================
    prs.save(OUTPUT_PATH)
    print(f"\nPresentasi berhasil dibuat!")
    print(f"File: {OUTPUT_PATH}")
    print(f"Total slide: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
