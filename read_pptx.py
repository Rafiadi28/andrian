from pptx import Presentation
import sys

def read_presentation(file_path):
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(f"[{shape.shape_type}] {shape.text.strip()}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        read_presentation(sys.argv[1])
    else:
        print("Please provide a file path")
