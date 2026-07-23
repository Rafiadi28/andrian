import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation(r'd:\laragon\www\andrian\Presentasi_Sistem_Kredit_BPR_Wonosobo.pptx')

# ============================================================
# SLIDE 5 (index 4) - Tambah keterangan Kepatuhan & maks 50jt
# ============================================================
slide5 = prs.slides[4]

dark_blue = RGBColor(29, 56, 88)
gold = RGBColor(196, 154, 69)
white = RGBColor(255, 255, 255)
red = RGBColor(192, 57, 43)

# 1. Add "Keterangan Plafond" note below jenis kredit (after Shape 10, below the boxes)
# Jenis kredit boxes end at top=2286000 + height=1005840 = 3291840
# "6 TAHAP" title is at top=3566160
# We have space between 3291840 and 3566160

note_box = slide5.shapes.add_textbox(Emu(731520), Emu(3291840), Emu(10698480), Emu(274320))
tf = note_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "\u26a0\ufe0f PENTING: Jenis kredit di atas terbatas untuk maksimal plafond kredit Rp 50.000.000 (lima puluh juta rupiah)"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = red

# 2. Add Kepatuhan status info - add new section below the 6 tahap
# The 6 tahap cards end at top=5257800+1097280 = 6355080
# Add below that

kepatuhan_title = slide5.shapes.add_textbox(Emu(731520), Emu(6492240), Emu(10058400), Emu(365760))
tf2 = kepatuhan_title.text_frame
p2 = tf2.paragraphs[0]
p2.text = "STATUS PENGAJUAN KREDIT (TERMASUK KEPATUHAN):"
p2.font.size = Pt(14)
p2.font.bold = True
p2.font.color.rgb = dark_blue

# Status boxes - in a row
statuses = [
    ("DRAFT", "Analis buat pengajuan baru", RGBColor(149, 165, 166)),
    ("DIAJUKAN", "Dikirim ke rantai approval", RGBColor(41, 128, 185)),
    ("KEPATUHAN", "Review kepatuhan & SOP", RGBColor(230, 126, 34)),
    ("PROSES", "Ditinjau pejabat berwenang", RGBColor(39, 174, 96)),
    ("DISETUJUI", "Kredit disetujui", RGBColor(46, 204, 113)),
    ("DITOLAK", "Kredit ditolak", RGBColor(231, 76, 60)),
]

status_y = 6858000
status_w = Emu(1737360)
status_h = Emu(640080)
status_x_start = 731520
status_gap = 137160

for idx, (label, desc, color) in enumerate(statuses):
    x = Emu(status_x_start + idx * (1737360 + status_gap))
    
    box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(status_y), status_w, status_h)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = label
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = white
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(9)
    p2.font.bold = False
    p2.font.color.rgb = white


# ============================================================
# SLIDE 7 (index 6) - Perlebar kotak kiri agar teks terbaca
# ============================================================
slide7 = prs.slides[6]

# The left-side cards (shapes 3-26) are currently:
#   left=457200, width=5303520
# The text boxes inside have left=594360 or 1051560
# The description text boxes have width=5029200 and height=182880 (terlalu kecil!)
# 
# We need to:
# 1. Make description text boxes taller (from 182880 to at least 457200)
# 2. Ensure proper word wrap

# Fix description text boxes (shapes 6, 10, 14, 18, 22, 26)
desc_shape_indices = [6, 10, 14, 18, 22, 26]

for idx in desc_shape_indices:
    shape = slide7.shapes[idx]
    # Make height bigger so text is fully readable
    shape.height = Emu(365760)  # was 182880, now doubled
    shape.width = Emu(5029200)
    if shape.has_text_frame:
        shape.text_frame.word_wrap = True
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is None or r.font.size < Pt(11):
                    r.font.size = Pt(11)

# Also make the background card shapes taller to accommodate
# Background cards: shapes 3, 7, 11, 15, 19, 23
bg_shape_indices = [3, 7, 11, 15, 19, 23]

for idx in bg_shape_indices:
    shape = slide7.shapes[idx]
    shape.height = Emu(914400)  # was 777240, increase

# Title text boxes (shapes 5, 9, 13, 17, 21, 25) - make font bigger
title_shape_indices = [5, 9, 13, 17, 21, 25]

for idx in title_shape_indices:
    shape = slide7.shapes[idx]
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(14)
                r.font.bold = True


prs.save(r'd:\laragon\www\andrian\Presentasi_Sistem_Kredit_BPR_Wonosobo.pptx')
print('Slide 5 dan Slide 7 berhasil diperbarui.')
