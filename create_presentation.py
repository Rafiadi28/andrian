from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()

# Helper functions

def add_text_slide(title_text, body_lines, back_slide=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for idx, line in enumerate(body_lines):
        if idx == 0:
            p = tf.paragraphs[0]
            p.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(18)
    if back_slide is not None:
        add_action_button(slide, Inches(7.2), Inches(5.8), Inches(2.0), Inches(0.6), 'Kembali ke Menu', back_slide)
    return slide


def add_action_button(slide, left, top, width, height, text, target_slide):
    btn = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height
    )
    btn.fill.solid()
    btn.fill.fore_color.rgb = RGBColor(49, 130, 206)
    btn.line.color.rgb = RGBColor(30, 64, 175)
    text_frame = btn.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    btn.click_action.target_slide = target_slide
    return btn


def create_usecase_slide(menu_slide):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.text = 'Use Case Sistem Kredit'
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    actor_left = Inches(0.3)
    actor_top = Inches(1.5)
    actor_height = Inches(0.8)
    actor_width = Inches(1.5)
    actors = ['Analis', 'Kepatuhan', 'Kabag Kredit', 'Direktur Utama']

    for idx, actor in enumerate(actors):
        actor_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            actor_left,
            actor_top + idx * (actor_height + Inches(0.15)),
            actor_width,
            actor_height
        )
        actor_shape.text = actor
        actor_tf = actor_shape.text_frame
        actor_tf.paragraphs[0].font.size = Pt(14)
        actor_tf.paragraphs[0].font.bold = True
        actor_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        actor_shape.fill.solid()
        actor_shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
        actor_shape.line.color.rgb = RGBColor(15, 23, 42)

    usecase_left = Inches(3.0)
    usecase_top = Inches(1.5)
    usecase_width = Inches(4.0)
    usecase_height = Inches(1.0)
    usecases = [
        'Input Data Pengajuan',
        'Verifikasi Kepatuhan',
        'Review Approval',
        'Approval Final'
    ]

    for idx, case in enumerate(usecases):
        oval = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            usecase_left,
            usecase_top + idx * (usecase_height + Inches(0.2)),
            usecase_width,
            usecase_height
        )
        oval.text = case
        oval_tf = oval.text_frame
        oval_tf.paragraphs[0].font.size = Pt(14)
        oval_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor(222, 234, 246)
        oval.line.color.rgb = RGBColor(49, 130, 206)

    relations = [
        (0, 0),
        (1, 1),
        (2, 2),
        (3, 3)
    ]
    for actor_idx, usecase_idx in relations:
        x1 = actor_left + actor_width
        y1 = actor_top + actor_idx * (actor_height + Inches(0.15)) + actor_height / 2
        x2 = usecase_left
        y2 = usecase_top + usecase_idx * (usecase_height + Inches(0.2)) + usecase_height / 2
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x1,
            y1,
            x2,
            y2
        )
        connector.line.width = Pt(2)
        connector.line.color.rgb = RGBColor(30, 64, 175)

    note = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9.0), Inches(0.8))
    note_tf = note.text_frame
    note_tf.text = 'Use case ini menunjukkan peran utama dalam sistem dan fungsi utama yang dilakukan setiap aktor selama proses kredit.'
    note_tf.paragraphs[0].font.size = Pt(16)
    note_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    add_action_button(slide, Inches(7.2), Inches(5.8), Inches(2.0), Inches(0.6), 'Kembali ke Menu', menu_slide)
    return slide


def create_flowchart_slide(menu_slide):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.text = 'Flowchart Proses Pengajuan'
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    steps = [
        'Mulai',
        'Analis Input Data',
        'Compliance Check',
        'Approval Review',
        'Nilai Kredit >= 500J?',
        'Final Approval',
        'Selesai'
    ]
    left = Inches(4.0)
    top = Inches(1.5)
    width = Inches(3.5)
    height = Inches(0.7)

    for idx, text in enumerate(steps):
        shape_type = MSO_AUTO_SHAPE_TYPE.OVAL if idx in (0, 6) else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type,
            left,
            top + idx * (height + Inches(0.25)),
            width,
            height
        )
        shape.text = text
        shape.text_frame.paragraphs[0].font.size = Pt(14)
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(236, 247, 255)
        shape.line.color.rgb = RGBColor(15, 23, 42)

    for idx in range(len(steps) - 1):
        x1 = left + width / 2
        y1 = top + idx * (height + Inches(0.25)) + height
        x2 = x1
        y2 = top + (idx + 1) * (height + Inches(0.25))
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x1,
            y1,
            x2,
            y2
        )
        connector.line.width = Pt(2)
        connector.line.color.rgb = RGBColor(15, 23, 42)
        connector.line.end_arrowhead = True

    note = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9.0), Inches(0.8))
    note_tf = note.text_frame
    note_tf.text = 'Flowchart ini menampilkan rangkaian proses dari input awal hingga final approval beserta keputusan threshold kredit.'
    note_tf.paragraphs[0].font.size = Pt(16)
    note_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    add_action_button(slide, Inches(7.2), Inches(5.8), Inches(2.0), Inches(0.6), 'Kembali ke Menu', menu_slide)
    return slide


def create_diagram_slide(menu_slide):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9.0), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.text = 'Diagram Alur Sistem'
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    nodes = [
        'Analis',
        'Kasubag Analis',
        'Kabag Kredit',
        'Kadiv Bisnis',
        'Direktur Utama'
    ]
    left = Inches(0.4)
    top = Inches(1.5)
    width = Inches(1.8)
    height = Inches(0.85)

    for idx, label in enumerate(nodes):
        node = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left + idx * (width + Inches(0.22)),
            top,
            width,
            height
        )
        node.text = label
        tf = node.text_frame
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        node.fill.solid()
        node.fill.fore_color.rgb = RGBColor(235, 248, 255)
        node.line.color.rgb = RGBColor(13, 60, 97)

    for idx in range(len(nodes) - 1):
        x1 = left + (idx + 1) * (width + Inches(0.22)) - Inches(0.03)
        y1 = top + height / 2
        x2 = x1 + Inches(0.22)
        y2 = y1
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x1,
            y1,
            x2,
            y2
        )
        connector.line.width = Pt(3)
        connector.line.color.rgb = RGBColor(13, 60, 97)
        connector.line.end_arrowhead = True

    note = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9.0), Inches(1.2))
    note_tf = note.text_frame
    note_tf.text = 'Pengajuan kredit diproses secara berjenjang dengan auto-skip jika jabatan tidak aktif, dan routing nominal < 500 juta berhenti di Kadiv Bisnis.'
    note_tf.paragraphs[0].font.size = Pt(16)
    note_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    add_action_button(slide, Inches(7.2), Inches(5.8), Inches(2.0), Inches(0.6), 'Kembali ke Menu', menu_slide)
    return slide


# Slide 1: Judul
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
slide_title.shapes.title.text = 'Alur Sistem Aplikasi Kredit'
subtitle = slide_title.placeholders[1]
subtitle.text = 'Presentasi interaktif untuk pimpinan\nSistem approval kredit dengan compliance terintegrasi'

# Slide 2: Menu interaktif
slide_menu = prs.slides.add_slide(prs.slide_layouts[1])
slide_menu.shapes.title.text = 'Menu Interaktif'
body = slide_menu.shapes.placeholders[1]
tf = body.text_frame
tf.text = 'Klik salah satu bagian untuk langsung menuju topik:'
for item in [
    '1. Inti Sistem',
    '2. Alur Pengajuan',
    '3. Approval Chain',
    '4. Revisi & Tolak',
    '5. Nilai Bisnis',
    '6. Use Case Sistem',
    '7. Flowchart Proses',
    '8. Diagram Alur Sistem'
]:
    p = tf.add_paragraph()
    p.text = item
    p.level = 0
    p.font.size = Pt(18)

button_texts = [
    ('Inti Sistem', 0),
    ('Alur Pengajuan', 1),
    ('Approval Chain', 2),
    ('Revisi & Tolak', 3),
    ('Nilai Bisnis', 4),
    ('Use Case Sistem', 5),
    ('Flowchart Proses', 6),
    ('Diagram Alur', 7)
]

# Create placeholder slides first for linking
section_slides = []
section_slides.append(add_text_slide('Inti Sistem', [
    'Sistem ini menangani seluruh proses pengajuan kredit: input analis, compliance, approval, dan cetak laporan.',
    'Memastikan data credit scoring, cashflow, agunan, dan kepatuhan tersimpan lengkap.',
    'Otomatis mengarahkan proses approval sesuai jabatan aktif dan batas nominal kredit.'
]))
section_slides.append(add_text_slide('Alur Pengajuan', [
    '1. Analis mengisi data debitur, usaha, cashflow, neraca, dan agunan.',
    '2. Data disimpan ke tabel pengajuan_kredit, analisa_neraca, assessment_kepatuhan, dan jaminan_.',
    '3. Status berubah menjadi diajukan / proses, dan sistem mengarahkan ke tahap approval berikutnya.',
    '4. Jika ada role tidak aktif, sistem auto-skip ke role berikutnya.'
]))
section_slides.append(add_text_slide('Approval Chain', [
    'Alur default: Analis → Kasubag Analis → Kabag Kredit → Kadiv Bisnis → Direktur Utama.',
    'Setiap role memiliki halaman kerja khusus untuk melihat pengajuan berdasarkan posisi_saat_ini.',
    'Approval dicatat di tabel approval_kredit sebagai audit trail.'
]))
section_slides.append(add_text_slide('Revisi & Tolak', [
    'Setuju: lanjut ke stage berikutnya atau selesai jika akhir chain.',
    'Revisi: kembali ke analis, catatan revisi disimpan, analis dapat perbaiki dan submit ulang.',
    'Tolak: status ditolak, analis dapat lihat alasan dan lakukan pengajuan ulang jika perlu.'
]))
section_slides.append(add_text_slide('Nilai Bisnis', [
    'Sistem mempercepat keputusan untuk kredit < 500 juta dengan berhenti di Kadiv Bisnis.',
    'Kredit ≥ 500 juta memerlukan approval final dari Direktur Utama.',
    'Model ini menjaga keseimbangan antara efisiensi dan kontrol risiko.'
]))
section_slides.append(create_usecase_slide(slide_menu))
section_slides.append(create_flowchart_slide(slide_menu))

# Add interactive buttons to menu after slides exist
for idx, (text, _) in enumerate(button_texts):
    target = section_slides[idx] if idx < len(section_slides) else None
    if target is not None:
        add_action_button(slide_menu, Inches(0.5 + (idx % 2) * 4.5), Inches(3.0 + (idx // 2) * 0.9), Inches(4.0), Inches(0.7), text, target)

# Add diagram slide
diagram_slide = create_diagram_slide(slide_menu)

# Add back button to title slide that links to menu
add_action_button(slide_title, Inches(7.2), Inches(5.8), Inches(2.0), Inches(0.6), 'Lihat Menu', slide_menu)

prs.save('Alur_Sistem_Aplikasi_Kredit_Interaktif_v3.pptx')
