import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu
prs = Presentation(r'd:\laragon\www\andrian\Presentasi_Sistem_Kredit_BPR_Wonosobo.pptx')

# Get slide dimensions
slide_w = prs.slide_width
slide_h = prs.slide_height
print(f"Slide dimensions: width={slide_w} ({Emu(slide_w).inches:.2f}in), height={slide_h} ({Emu(slide_h).inches:.2f}in)")

slide5 = prs.slides[4]
with open('slide5_full.txt', 'w', encoding='utf-8') as f:
    f.write(f'Slide: {slide_w}x{slide_h}\n')
    for i, shape in enumerate(slide5.shapes):
        right = shape.left + shape.width
        bottom = shape.top + shape.height
        offside = ""
        if right > slide_w or bottom > slide_h:
            offside = " *** OFFSIDE ***"
        f.write(f'Shape {i}: type={shape.shape_type}, L={shape.left}, T={shape.top}, W={shape.width}, H={shape.height}, R={right}, B={bottom}{offside}\n')
        if shape.has_text_frame:
            for pi, p in enumerate(shape.text_frame.paragraphs):
                f.write(f'  P{pi}: "{p.text}"\n')
        f.write('\n')
print("Done")
